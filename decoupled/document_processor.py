import boto3
import tabula
import base64
import pymupdf  # Fitz
import os
import warnings
from io import BytesIO
from PIL import Image
from botocore.exceptions import ClientError
from dotenv import load_dotenv

# LangChain & AWS Imports
from langchain_aws import ChatBedrockConverse
from langchain_core.messages import HumanMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import (
    PyMuPDFLoader,
    Docx2txtLoader,
    UnstructuredPowerPointLoader,
    TextLoader,
)

# Office Document Parsers
import docx
from docx.enum.shape import WD_INLINE_SHAPE
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---
# Import the new prompt loader
# ---
from decoupled.utils import load_prompt

# ---
# 1. BEDROCK & CLIENT INITIALIZATIONS
# ---

load_dotenv()
os.environ["AWS_BEARER_TOKEN_BEDROCK"] = os.getenv("BEDROCK_API", "")

try:
    boto_client = boto3.client(
        service_name="bedrock-runtime",
        region_name="us-east-1"
    )

    llm = ChatBedrockConverse(
        model="anthropic.claude-haiku-4-5-20251001-v1:0",
        client=boto_client,
        temperature=0.0
    )
    
    print("DocumentProcessor: Bedrock LLM initialized successfully.")

except Exception as e:
    print(f"Error initializing Boto3 client or LLM in DocumentProcessor: {e}")
    exit()

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=100,
    length_function=len
)

try:
    PROMPT_IMAGE_DESCRIPTION = load_prompt("image_description.txt")
except Exception as e:
    exit(1) # Exit if prompt loading fails

def create_output_directories(base_dir):
    """Creates the necessary output subdirectories."""
    # Main directories
    for dir_name in ["images", "text", "other"]:
        dir_path = os.path.join(base_dir, dir_name)
        os.makedirs(dir_path, exist_ok=True)
        print(f"Created directory: {dir_path}")
    
    # Subdirectories within 'other' for different content types (charts removed)
    other_subdirs = ["tables", "forms", "equations", "metadata", "annotations"]
    for subdir in other_subdirs:
        subdir_path = os.path.join(base_dir, "other", subdir)
        os.makedirs(subdir_path, exist_ok=True)
        print(f"Created subdirectory: {subdir_path}")


def convert_image_to_png_base64(image_bytes: bytes) -> str:
    """Converts any image format to PNG and returns base64 encoded string."""
    try:
        # Open image from bytes
        image = Image.open(BytesIO(image_bytes))
        
        print(f"   - Original image mode: {image.mode}, size: {image.size}")
        
        # Convert to RGB if necessary (handles RGBA, P mode, etc.)
        if image.mode in ('RGBA', 'LA', 'P'):
            # Create white background
            background = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'P':
                image = image.convert('RGBA')
            background.paste(image, mask=image.split()[-1] if image.mode in ('RGBA', 'LA') else None)
            image = background
        elif image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Save as PNG to bytes buffer
        buffer = BytesIO()
        image.save(buffer, format='PNG')
        png_size = buffer.tell()
        buffer.seek(0)
        
        # Encode to base64
        base64_string = base64.b64encode(buffer.read()).decode('utf-8')
        print(f"   - Converted to PNG, size: {png_size} bytes, base64 length: {len(base64_string)}")
        
        return base64_string
    except Exception as e:
        print(f"   - Error converting image to PNG: {e}")
        raise


def generate_image_description(image_bytes: bytes, context: str = "") -> str:
    """Generates a detailed description of an image using the Bedrock LLM."""

    prompt = PROMPT_IMAGE_DESCRIPTION
    
    try:
        print(f"   - Generating description for image {context}...")
        # Convert image to PNG base64
        base64_image = convert_image_to_png_base64(image_bytes)
        
        messages = [
            HumanMessage(content=[
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}},
            ])
        ]
        
        print(f"   - Calling Bedrock API...")
        # 'llm' is now defined at the top of this file
        response = llm.invoke(messages)
        description = response.content
        print(f"   - ✓ Description generated successfully ({len(description)} chars)")
        return description
    except Exception as e:
        error_msg = f"Error generating image description: {str(e)}"
        print(f"   - ✗ {error_msg}")
        return f"Error: {error_msg}"


# ... (rest of the file remains unchanged) ...

def extract_document_metadata(filepath, base_dir, items):
    """Extract document metadata and save to other/metadata/"""
    file_basename = os.path.basename(filepath)
    metadata = {}
    
    try:
        if filepath.lower().endswith('.pdf'):
            doc = pymupdf.open(filepath)
            metadata = doc.metadata
            doc.close()
        elif filepath.lower().endswith('.docx'):
            doc = docx.Document(filepath)
            core_props = doc.core_properties
            metadata = {
                'title': core_props.title,
                'author': core_props.author,
                'subject': core_props.subject,
                'keywords': core_props.keywords,
                'created': str(core_props.created),
                'modified': str(core_props.modified),
                'last_modified_by': core_props.last_modified_by,
                'revision': core_props.revision,
                'category': core_props.category,
                'comments': core_props.comments
            }
        elif filepath.lower().endswith('.pptx'):
            prs = Presentation(filepath)
            core_props = prs.core_properties
            metadata = {
                'title': core_props.title,
                'author': core_props.author,
                'subject': core_props.subject,
                'keywords': core_props.keywords,
                'created': str(core_props.created),
                'modified': str(core_props.modified),
                'last_modified_by': core_props.last_modified_by,
                'revision': core_props.revision,
                'category': core_props.category,
                'comments': core_props.comments
            }
        
        # Save metadata
        if metadata:
            metadata_text = "\n".join([f"{k}: {v}" for k, v in metadata.items() if v])
            path = f"{base_dir}/other/metadata/{file_basename}_metadata.txt"
            with open(path, 'w', encoding='utf-8') as f:
                f.write(metadata_text)
            items.append({"page": 0, "type": "metadata", "text": metadata_text, "path": path})
            print(f"   ✓ Extracted metadata: {len(metadata)} properties")
            
    except Exception as e:
        print(f"   - Error extracting metadata: {e}")


def extract_pdf_elements(filepath, base_dir, items):
    """Extracts tables and images from PDF files."""
    file_basename = os.path.basename(filepath)
    doc = pymupdf.open(filepath)
    
    print(f"\nExtracting from {file_basename} ({len(doc)} pages)...")
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        print(f"\n--- Page {page_num + 1} ---")
        
        # Table extraction
        try:
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                tables = tabula.read_pdf(filepath, pages=page_num + 1, multiple_tables=True, silent=True)
                if tables:
                    for i, table_df in enumerate(tables):
                        if not table_df.empty:  # Check if table has content
                            table_text = table_df.to_string()
                            path = f"{base_dir}/other/tables/{file_basename}_p{page_num}_t{i}.txt"
                            with open(path, 'w', encoding='utf-8') as f:
                                f.write(table_text)
                            print(f"   ✓ Extracted table {i} ({len(table_text)} chars)")
                            items.append({"page": page_num, "type": "table", "text": table_text, "path": path})
        except Exception as e:
            print(f"   - No tables found or extraction failed")
        
        # Image extraction and description
        images = page.get_images(full=True)
        print(f"   Found {len(images)} image(s) on page {page_num + 1}")
        
        for img_idx, img in enumerate(images):
            xref = img[0]
            print(f"\n   Processing image {img_idx + 1}/{len(images)}...")
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                print(f"   - Image format: {image_ext}, size: {len(image_bytes)} bytes")
                
                # Save original image
                image_path = f"{base_dir}/images/{file_basename}_p{page_num}_img{img_idx}.{image_ext}"
                with open(image_path, 'wb') as img_file:
                    img_file.write(image_bytes)
                print(f"   - ✓ Saved original image to: {image_path}")
                
                # Generate description using raw bytes
                description = generate_image_description(image_bytes, f"page {page_num + 1}, img {img_idx + 1}")
                
                if description and not description.startswith("Error:"):
                    desc_path = f"{base_dir}/text/{file_basename}_p{page_num}_img{img_idx}_desc.txt"
                    with open(desc_path, 'w', encoding='utf-8') as f:
                        f.write(description)
                    print(f"   - ✓ Saved description to: {desc_path}")
                    items.append({"page": page_num, "type": "image_description", "text": description, "path": desc_path})
                else:
                    print(f"   - ✗ Failed to generate description")
                    
            except Exception as e:
                print(f"   - ✗ Could not process image {img_idx}: {e}")
                continue
    
    doc.close()
    print(f"\nFinished extracting from {file_basename}")


def extract_docx_images(filepath, base_dir, items):
    """Extracts images from DOCX files and generates descriptions."""
    file_basename = os.path.basename(filepath)
    
    try:
        doc = docx.Document(filepath)
        
        print(f"\nExtracting images from {file_basename}...")
        print(f"Found {len(doc.inline_shapes)} inline shapes")
        
        image_count = 0
        for shape_idx, shape in enumerate(doc.inline_shapes):
            if shape.type == WD_INLINE_SHAPE.PICTURE:
                print(f"\nProcessing DOCX image {image_count + 1}...")
                try:
                    image_bytes_ref = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                    image_part = doc.part.related_parts[image_bytes_ref]
                    image_bytes = image_part.blob
                    
                    print(f"   - Image size: {len(image_bytes)} bytes")
                    
                    # Save original image
                    ext = image_part.content_type.split('/')[-1]
                    image_path = f"{base_dir}/images/{file_basename}_img{image_count}.{ext}"
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    print(f"   - ✓ Saved original image to: {image_path}")
                    
                    description = generate_image_description(image_bytes, f"DOCX img {image_count + 1}")
                    
                    if description and not description.startswith("Error:"):
                        desc_path = f"{base_dir}/text/{file_basename}_img{image_count}_desc.txt"
                        with open(desc_path, 'w', encoding='utf-8') as f:
                            f.write(description)
                        print(f"   - ✓ Saved description to: {desc_path}")
                        items.append({"page": 0, "type": "image_description", "text": description, "path": desc_path})
                        image_count += 1
                except Exception as e:
                    print(f"   - ✗ Could not process DOCX image: {e}")
                    continue
        
        print(f"\nExtracted {image_count} images from {file_basename}")
    except Exception as e:
        print(f"Error processing DOCX file: {e}")


def extract_pptx_images(filepath, base_dir, items):
    """Extracts images from PPTX files and generates descriptions."""
    file_basename = os.path.basename(filepath)
    
    try:
        prs = Presentation(filepath)
        print(f"\nExtracting images from {file_basename}...")
        print(f"Found {len(prs.slides)} slides")
        
        total_images = 0
        for slide_num, slide in enumerate(prs.slides):
            print(f"\n--- Slide {slide_num + 1} ---")
            slide_images = 0
            
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    print(f"   Processing PPTX image {shape_idx}...")
                    try:
                        image_bytes = shape.image.blob
                        ext = shape.image.content_type.split('/')[-1]
                        
                        print(f"   - Image format: {ext}, size: {len(image_bytes)} bytes")
                        
                        # Save original image
                        image_path = f"{base_dir}/images/{file_basename}_s{slide_num}_img{shape_idx}.{ext}"
                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_bytes)
                        print(f"   - ✓ Saved original image to: {image_path}")
                        
                        description = generate_image_description(image_bytes, f"slide {slide_num + 1}, img {shape_idx}")
                        
                        if description and not description.startswith("Error:"):
                            desc_path = f"{base_dir}/text/{file_basename}_s{slide_num}_img{shape_idx}_desc.txt"
                            with open(desc_path, 'w', encoding='utf-8') as f:
                                f.write(description)
                            print(f"   - ✓ Saved description to: {desc_path}")
                            items.append({"page": slide_num, "type": "image_description", "text": description, "path": desc_path})
                            slide_images += 1
                            total_images += 1
                    except Exception as e:
                        print(f"   - ✗ Could not process PPTX image: {e}")
                        continue
            
            if slide_images == 0:
                print(f"   No images found on this slide")
        
        print(f"\nExtracted {total_images} images from {file_basename}")
    except Exception as e:
        print(f"Error processing PPTX file: {e}")


def load_and_process_document(filepath):
    """
    Loads a document, chunks text, and calls specialized functions
    for non-text elements like images and tables.
    """
    if not os.path.exists(filepath):
        print(f"Error: File not found at '{filepath}'")
        return []

    file_root, file_ext = os.path.splitext(os.path.basename(filepath))
    file_ext = file_ext.lower()
    output_base_dir = f"output_{file_root}"
    
    print(f"\n{'='*60}")
    print(f"Creating output directory: {output_base_dir}")
    print(f"{'='*60}")
    create_output_directories(output_base_dir)

    loader_map = {
        '.pdf': PyMuPDFLoader,
        '.docx': Docx2txtLoader,
        '.pptx': UnstructuredPowerPointLoader,
        '.txt': TextLoader,
    }
    
    loader_class = loader_map.get(file_ext)
    if not loader_class:
        print(f"Unsupported file type: {file_ext}")
        return []

    # Load and Chunk Text Content
    print(f"\n{'='*60}")
    print(f"Loading text with {loader_class.__name__}...")
    print(f"{'='*60}")
    try:
        loader = loader_class(filepath)
        docs = loader.load()
        # 'text_splitter' is now defined at the top of this file
        chunks = text_splitter.split_documents(docs)
        processed_items = []
        
        print(f"Generated {len(chunks)} text chunks")
        
        for i, chunk in enumerate(chunks):
            page_num = chunk.metadata.get('page', 0)
            path = f"{output_base_dir}/text/{file_root}_text_chunk_{i}.txt"
            with open(path, 'w', encoding='utf-8') as f:
                f.write(chunk.page_content)
            processed_items.append({
                "page": page_num,
                "type": "text",
                "text": chunk.page_content,
                "path": path
            })
        print(f"✓ Saved {len(chunks)} text chunks")
        
    except Exception as e:
        print(f"Error loading document: {e}")
        return []
    
    # Call specialized extractors for non-text elements
    print(f"\n{'='*60}")
    print(f"Extracting non-text elements...")
    print(f"{'='*60}")
    
    if file_ext == '.pdf':
        print(f"     Extracting comprehensive content from PDF...")
        # Extract metadata first
        extract_document_metadata(filepath, output_base_dir, processed_items)
        # Extract tables and images
        extract_pdf_elements(filepath, output_base_dir, processed_items)
    elif file_ext == '.docx':
        print(f"     Extracting comprehensive content from DOCX...")
        # Extract metadata first
        extract_document_metadata(filepath, output_base_dir, processed_items)
        # Extract images
        extract_docx_images(filepath, output_base_dir, processed_items)
    elif file_ext == '.pptx':
        print(f"     Extracting comprehensive content from PPTX...")
        # Extract metadata first
        extract_document_metadata(filepath, output_base_dir, processed_items)
        # Extract images
        extract_pptx_images(filepath, output_base_dir, processed_items)
        
    return processed_items